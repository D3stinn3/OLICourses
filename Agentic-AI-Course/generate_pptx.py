"""
Generate a PowerPoint presentation from the Agentic AI Course HTML presentation.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors from the HTML theme
DARK_BG = RGBColor(0x0A, 0x0A, 0x1A)
DARK_BG2 = RGBColor(0x0F, 0x0F, 0x2E)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
PINK = RGBColor(0xF4, 0x72, 0xB6)
LIGHT_TEXT = RGBColor(0xCB, 0xD5, 0xE1)
MUTED_TEXT = RGBColor(0x94, 0xA3, 0xB8)
DIM_TEXT = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xE0, 0xE0, 0xE0)
HIGHLIGHT_BG = RGBColor(0x1E, 0x3A, 0x5F)
TABLE_HEADER_BG = RGBColor(0x15, 0x2A, 0x4A)
TABLE_ROW_BG = RGBColor(0x0F, 0x14, 0x2E)
TABLE_ALT_BG = RGBColor(0x12, 0x1A, 0x35)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, font_size=18, color=LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_paragraph(tf, text, font_size=18, color=LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                   space_before=Pt(0), space_after=Pt(6)):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_rich_paragraph(tf, segments, font_size=18, alignment=PP_ALIGN.LEFT,
                        space_before=Pt(0), space_after=Pt(6)):
    """Add a paragraph with mixed bold/normal/italic segments.
    segments: list of (text, color, bold, italic) tuples
    """
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    for text, color, bold, italic in segments:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
    return p


def add_bullet(tf, text, font_size=20, color=LIGHT_TEXT, bullet_color=PURPLE,
               space_before=Pt(4), space_after=Pt(6)):
    """Add a bullet point with the triangle marker."""
    p = tf.add_paragraph()
    p.space_before = space_before
    p.space_after = space_after

    bullet_run = p.add_run()
    bullet_run.text = "\u25B8  "
    bullet_run.font.size = Pt(font_size)
    bullet_run.font.color.rgb = bullet_color
    bullet_run.font.bold = True

    text_run = p.add_run()
    text_run.text = text
    text_run.font.size = Pt(font_size)
    text_run.font.color.rgb = color
    return p


def add_rich_bullet(tf, segments, font_size=20, bullet_color=PURPLE,
                     space_before=Pt(4), space_after=Pt(6)):
    """Add a bullet with mixed formatting segments."""
    p = tf.add_paragraph()
    p.space_before = space_before
    p.space_after = space_after

    bullet_run = p.add_run()
    bullet_run.text = "\u25B8  "
    bullet_run.font.size = Pt(font_size)
    bullet_run.font.color.rgb = bullet_color
    bullet_run.font.bold = True

    for text, color, bold, italic in segments:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
    return p


def add_highlight_box(slide, left, top, width, height, text_segments, font_size=18):
    """Add a highlight box (blue left border with background)."""
    # Background rectangle
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x12, 0x25, 0x40)
    shape.line.fill.background()

    # Blue left accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    # Text
    txBox = add_textbox(slide, left + Inches(0.3), top + Inches(0.15),
                        width - Inches(0.5), height - Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for segments in text_segments:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        for text, color, bold, italic in segments:
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.italic = italic


def add_diagram_flow(slide, items, top, left_start=None, box_width=Inches(2.2),
                      box_height=Inches(1.0), gap=Inches(0.3), arrow_width=Inches(0.5)):
    """Add a horizontal flow diagram with boxes and arrows."""
    total_width = len(items) * box_width + (len(items) - 1) * (arrow_width + gap * 2)
    if left_start is None:
        left_start = (SLIDE_W - total_width) // 2

    x = left_start
    for i, (title, subtitle) in enumerate(items):
        # Box
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_width, box_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x1A, 0x15, 0x35)
        shape.line.color.rgb = PURPLE
        shape.line.width = Pt(2)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = title
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = WHITE

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(12)
        r2.font.color.rgb = MUTED_TEXT

        x += box_width + gap

        # Arrow (except after last box)
        if i < len(items) - 1:
            arrow_tb = add_textbox(slide, x, top + Inches(0.25), arrow_width, Inches(0.5))
            atf = arrow_tb.text_frame
            atf.paragraphs[0].alignment = PP_ALIGN.CENTER
            ar = atf.paragraphs[0].add_run()
            ar.text = "\u2192"
            ar.font.size = Pt(28)
            ar.font.color.rgb = BLUE
            x += arrow_width + gap


def add_slide_title(slide, text, left=Inches(0.8), top=Inches(0.5), width=Inches(11.5)):
    """Add a styled section title with underline."""
    txBox = add_textbox(slide, left, top, width, Inches(0.8))
    tf = txBox.text_frame
    set_text(tf, text, font_size=36, color=BLUE, bold=True)

    # Underline bar
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top + Inches(0.75), Inches(10), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = HIGHLIGHT_BG
    line.line.fill.background()


def add_table(slide, headers, rows, left=Inches(0.8), top=Inches(2.0), width=Inches(11.5),
              col_widths=None):
    """Add a styled table."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, Inches(0.5 * num_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(16)
                run.font.color.rgb = BLUE
                run.font.bold = True

    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            bg = TABLE_ROW_BG if r % 2 == 0 else TABLE_ALT_BG
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(14)
                    run.font.color.rgb = LIGHT_TEXT

    # Set first row as header
    tbl = table._tbl
    tblPr = tbl.tblPr if tbl.tblPr is not None else tbl._add_tblPr()
    tblPr.set('firstRow', '1')


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide1, RGBColor(0x0A, 0x0A, 0x2E))

# Title
txBox = add_textbox(slide1, Inches(1), Inches(2.0), Inches(11.3), Inches(1.5))
tf = txBox.text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
run = tf.paragraphs[0].add_run()
run.text = "Agentic AI"
run.font.size = Pt(60)
run.font.bold = True
run.font.color.rgb = BLUE

# Subtitle
txBox2 = add_textbox(slide1, Inches(1), Inches(3.5), Inches(11.3), Inches(0.8))
tf2 = txBox2.text_frame
set_text(tf2, "Building Autonomous Intelligent Systems", font_size=28, color=MUTED_TEXT,
         alignment=PP_ALIGN.CENTER)

# Meta
txBox3 = add_textbox(slide1, Inches(1), Inches(4.5), Inches(11.3), Inches(0.6))
tf3 = txBox3.text_frame
set_text(tf3, "OLI-AAI-101 | OpenClaw Learning Initiative", font_size=16, color=DIM_TEXT,
         alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: What is Agentic AI?
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, DARK_BG2)
add_slide_title(slide2, "What is Agentic AI?")

# Description
txBox = add_textbox(slide2, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
add_rich_paragraph(tf, [
    ("Agentic AI refers to AI systems that can ", LIGHT_TEXT, False, False),
    ("independently", LIGHT_TEXT, True, False),
    (" perceive their environment, make decisions, take actions, and learn from outcomes \u2014 with minimal human guidance.", LIGHT_TEXT, False, False),
], font_size=20)

# Highlight box
add_highlight_box(slide2, Inches(0.8), Inches(2.6), Inches(11.5), Inches(1.1), [
    [("Key Idea: ", BLUE, True, False),
     ("An AI agent doesn't just answer questions \u2014 it ", LIGHT_TEXT, False, False),
     ("does things", LIGHT_TEXT, False, True),
     (". It reasons about what to do, uses tools, and works toward goals on its own.", LIGHT_TEXT, False, False)]
])

# Bullets
txBox2 = add_textbox(slide2, Inches(0.8), Inches(4.0), Inches(11.5), Inches(2.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_bullet(tf2, "Goes beyond simple chatbots and autocomplete", font_size=20)
add_bullet(tf2, "Can break complex tasks into steps and execute them", font_size=20)
add_bullet(tf2, "Interacts with the real world through tools and APIs", font_size=20)

# ============================================================
# SLIDE 3: Traditional AI vs. Agentic AI
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, DARK_BG2)
add_slide_title(slide3, "Traditional AI vs. Agentic AI")

# Left column: Traditional AI
left_x = Inches(0.8)
col_w = Inches(5.5)

txL = add_textbox(slide3, left_x, Inches(1.6), col_w, Inches(0.5))
set_text(txL.text_frame, "Traditional AI / Chatbots", font_size=24, color=PURPLE, bold=True)

txL2 = add_textbox(slide3, left_x, Inches(2.2), col_w, Inches(4.0))
tfL = txL2.text_frame
tfL.word_wrap = True
for item in [
    "Responds only when prompted",
    "Single turn: question in, answer out",
    "No memory across interactions",
    "Cannot take actions in the world",
    "Human does all the planning",
]:
    add_bullet(tfL, item, font_size=18)

# Right column: Agentic AI
right_x = Inches(7.0)

txR = add_textbox(slide3, right_x, Inches(1.6), col_w, Inches(0.5))
set_text(txR.text_frame, "Agentic AI", font_size=24, color=PURPLE, bold=True)

txR2 = add_textbox(slide3, right_x, Inches(2.2), col_w, Inches(4.0))
tfR = txR2.text_frame
tfR.word_wrap = True
for item in [
    "Proactively works toward goals",
    "Multi-step reasoning and execution",
    "Maintains memory and context",
    "Uses tools (search, code, APIs)",
    "Plans, acts, and self-corrects",
]:
    add_bullet(tfR, item, font_size=18)

# ============================================================
# SLIDE 4: Core Components of an AI Agent
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, DARK_BG2)
add_slide_title(slide4, "Core Components of an AI Agent")

add_diagram_flow(slide4, [
    ("Perception", "Observe inputs"),
    ("Reasoning", "Think & plan"),
    ("Action", "Use tools & act"),
    ("Learning", "Reflect & improve"),
], top=Inches(1.8))

txBox = add_textbox(slide4, Inches(0.8), Inches(3.3), Inches(11.5), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True
add_rich_bullet(tf, [("Perception: ", BLUE, True, False), ("The agent receives input \u2014 a user request, sensor data, or system event", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf, [("Reasoning: ", BLUE, True, False), ("An LLM (like Claude or GPT) thinks through what to do next", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf, [("Action: ", BLUE, True, False), ("The agent calls tools, writes code, searches the web, or sends messages", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf, [("Learning: ", BLUE, True, False), ("The agent reflects on results and adjusts its approach", LIGHT_TEXT, False, False)], font_size=18)

# ============================================================
# SLIDE 5: LLMs: The Brain of the Agent
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, DARK_BG2)
add_slide_title(slide5, "LLMs: The Brain of the Agent")

txBox = add_textbox(slide5, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
add_rich_paragraph(tf, [
    ("Large Language Models (LLMs) like Claude, GPT, and Gemini serve as the ", LIGHT_TEXT, False, False),
    ("reasoning engine", LIGHT_TEXT, True, False),
    (" inside modern AI agents.", LIGHT_TEXT, False, False),
], font_size=20)

txBox2 = add_textbox(slide5, Inches(0.8), Inches(2.5), Inches(11.5), Inches(3.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_rich_bullet(tf2, [("Natural language understanding", BLUE, True, False), (" \u2014 interpret complex user requests", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf2, [("Planning", BLUE, True, False), (" \u2014 break big tasks into smaller steps", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf2, [("Decision making", BLUE, True, False), (" \u2014 choose which tool to use and when", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf2, [("Code generation", BLUE, True, False), (" \u2014 write and debug code on the fly", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf2, [("Self-reflection", BLUE, True, False), (" \u2014 evaluate their own outputs and retry if needed", LIGHT_TEXT, False, False)], font_size=18)

add_highlight_box(slide5, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.9), [
    [("The LLM doesn't just generate text \u2014 it ", LIGHT_TEXT, False, False),
     ("orchestrates", LIGHT_TEXT, False, True),
     (" an entire workflow of reasoning, tool use, and evaluation.", LIGHT_TEXT, False, False)]
])

# ============================================================
# SLIDE 6: The ReAct Pattern
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, DARK_BG2)
add_slide_title(slide6, "The ReAct Pattern (Reasoning + Acting)")

txBox = add_textbox(slide6, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7))
tf = txBox.text_frame
tf.word_wrap = True
add_rich_paragraph(tf, [
    ("ReAct is a foundational pattern where the agent alternates between ", LIGHT_TEXT, False, False),
    ("thinking", LIGHT_TEXT, True, False),
    (" and ", LIGHT_TEXT, False, False),
    ("doing", LIGHT_TEXT, True, False),
    (".", LIGHT_TEXT, False, False),
], font_size=20)

add_diagram_flow(slide6, [
    ("Thought", '"I need to find..."'),
    ("Action", 'search("query")'),
    ("Observation", "Results returned"),
    ("Thought", '"Now I know..."'),
], top=Inches(2.5))

txBox2 = add_textbox(slide6, Inches(0.8), Inches(4.0), Inches(11.5), Inches(3.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_rich_bullet(tf2, [("The agent ", LIGHT_TEXT, False, False), ("thinks out loud", LIGHT_TEXT, True, False), (" about what it needs to do", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf2, [("It then ", LIGHT_TEXT, False, False), ("takes an action", LIGHT_TEXT, True, False), (" (e.g., call a tool)", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf2, [("It ", LIGHT_TEXT, False, False), ("observes", LIGHT_TEXT, True, False), (" the result and decides the next step", LIGHT_TEXT, False, False)], font_size=18)
add_bullet(tf2, "This loop repeats until the task is complete", font_size=18)

# ============================================================
# SLIDE 7: Tool Use
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide7, DARK_BG2)
add_slide_title(slide7, "Tool Use: How Agents Interact with the World")

txBox = add_textbox(slide7, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7))
tf = txBox.text_frame
tf.word_wrap = True
add_rich_paragraph(tf, [
    ("Agents become truly useful when they can ", LIGHT_TEXT, False, False),
    ("use tools", LIGHT_TEXT, True, False),
    (" \u2014 external functions that extend their capabilities beyond text generation.", LIGHT_TEXT, False, False),
], font_size=20)

txBox2 = add_textbox(slide7, Inches(0.8), Inches(2.5), Inches(11.5), Inches(3.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_rich_bullet(tf2, [("Web Search", BLUE, True, False), (" \u2014 look up current information online", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf2, [("Code Execution", BLUE, True, False), (" \u2014 write and run Python, JavaScript, etc.", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf2, [("File Operations", BLUE, True, False), (" \u2014 read, write, and manage files", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf2, [("API Calls", BLUE, True, False), (" \u2014 interact with databases, services, and platforms", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf2, [("Calculators & Data Tools", BLUE, True, False), (" \u2014 perform precise computations", LIGHT_TEXT, False, False)], font_size=18)

add_highlight_box(slide7, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.0), [
    [("Example: ", BLUE, True, False),
     ('"What\'s the weather in Nairobi?" \u2014 The agent calls a weather API, gets live data, and presents the answer.', LIGHT_TEXT, False, False)]
])

# ============================================================
# SLIDE 8: Agent Memory
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide8, DARK_BG2)
add_slide_title(slide8, "Agent Memory: Remembering and Retrieving")

txBox = add_textbox(slide8, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7))
tf = txBox.text_frame
tf.word_wrap = True
add_rich_paragraph(tf, [
    ("Without memory, agents forget everything between interactions. Memory systems give agents ", LIGHT_TEXT, False, False),
    ("continuity", LIGHT_TEXT, True, False),
    (" and ", LIGHT_TEXT, False, False),
    ("context", LIGHT_TEXT, True, False),
    (".", LIGHT_TEXT, False, False),
], font_size=20)

# Left column: Short-Term Memory
txL = add_textbox(slide8, left_x, Inches(2.5), col_w, Inches(0.5))
set_text(txL.text_frame, "Short-Term Memory", font_size=22, color=PURPLE, bold=True)

txL2 = add_textbox(slide8, left_x, Inches(3.1), col_w, Inches(2.5))
tfL = txL2.text_frame
tfL.word_wrap = True
add_bullet(tfL, "Current conversation context", font_size=17)
add_bullet(tfL, "Working notes during a task", font_size=17)
add_bullet(tfL, "Limited by context window size", font_size=17)

# Right column: Long-Term Memory
txR = add_textbox(slide8, right_x, Inches(2.5), col_w, Inches(0.5))
set_text(txR.text_frame, "Long-Term Memory", font_size=22, color=PURPLE, bold=True)

txR2 = add_textbox(slide8, right_x, Inches(3.1), col_w, Inches(2.5))
tfR = txR2.text_frame
tfR.word_wrap = True
add_bullet(tfR, "Stored in vector databases", font_size=17)
add_bullet(tfR, "Persists across sessions", font_size=17)
add_bullet(tfR, "Retrieved via similarity search (RAG)", font_size=17)

add_highlight_box(slide8, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.0), [
    [("RAG (Retrieval-Augmented Generation): ", BLUE, True, False),
     ("The agent searches a knowledge base for relevant info and includes it in its reasoning.", LIGHT_TEXT, False, False)]
])

# ============================================================
# SLIDE 9: How Agents Plan
# ============================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide9, DARK_BG2)
add_slide_title(slide9, "How Agents Plan")

txBox = add_textbox(slide9, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7))
tf = txBox.text_frame
tf.word_wrap = True
add_rich_paragraph(tf, [
    ("For complex tasks, agents need to ", LIGHT_TEXT, False, False),
    ("plan ahead", LIGHT_TEXT, True, False),
    (" rather than just react step-by-step.", LIGHT_TEXT, False, False),
], font_size=20)

txBox2 = add_textbox(slide9, Inches(0.8), Inches(2.5), Inches(11.5), Inches(2.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_rich_bullet(tf2, [("Plan-and-Execute: ", BLUE, True, False), ("First create a full plan, then execute each step", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf2, [("Chain-of-Thought: ", BLUE, True, False), ("Reason through the problem step-by-step before acting", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf2, [("Tree-of-Thought: ", BLUE, True, False), ("Explore multiple reasoning paths and pick the best one", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf2, [("Reflexion: ", BLUE, True, False), ("After completing a task, reflect on what went wrong and improve", LIGHT_TEXT, False, False)], font_size=18)

add_highlight_box(slide9, Inches(0.8), Inches(4.8), Inches(11.5), Inches(1.4), [
    [("Example \u2014 Plan-and-Execute:", BLUE, True, False)],
    [("1. Research the topic \u2192 2. Outline the article \u2192 3. Write each section \u2192 4. Review and edit \u2192 5. Format and deliver", LIGHT_TEXT, False, False)],
], font_size=16)

# ============================================================
# SLIDE 10: Multi-Agent Systems
# ============================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide10, DARK_BG2)
add_slide_title(slide10, "Multi-Agent Systems")

txBox = add_textbox(slide10, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7))
tf = txBox.text_frame
tf.word_wrap = True
add_rich_paragraph(tf, [
    ("Instead of one agent doing everything, ", LIGHT_TEXT, False, False),
    ("multiple specialized agents", LIGHT_TEXT, True, False),
    (" can collaborate \u2014 each with its own role and expertise.", LIGHT_TEXT, False, False),
], font_size=20)

# Orchestrator diagram
orch_x = Inches(2.5)
orch_top = Inches(2.5)
box_w = Inches(2.2)
box_h = Inches(0.8)

# Orchestrator box
shape = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, orch_x, orch_top + Inches(0.8), box_w, box_h)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x15, 0x35)
shape.line.color.rgb = PINK
shape.line.width = Pt(2)
stf = shape.text_frame
stf.paragraphs[0].alignment = PP_ALIGN.CENTER
sr = stf.paragraphs[0].add_run()
sr.text = "Orchestrator\nAssigns tasks"
sr.font.size = Pt(14)
sr.font.color.rgb = WHITE
sr.font.bold = True

# Agent boxes
agents = [("Researcher", "Finds info"), ("Writer", "Drafts content"), ("Reviewer", "Checks quality")]
agent_x = Inches(6.5)
for i, (name, desc) in enumerate(agents):
    y = orch_top + Inches(i * 1.0)
    # Arrow
    atb = add_textbox(slide10, orch_x + box_w + Inches(0.2), y + Inches(0.15), Inches(1.5), Inches(0.5))
    atf = atb.text_frame
    atf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ar = atf.paragraphs[0].add_run()
    ar.text = "\u2192"
    ar.font.size = Pt(24)
    ar.font.color.rgb = BLUE

    # Box
    s = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, agent_x, y, box_w, box_h)
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0x1A, 0x15, 0x35)
    s.line.color.rgb = PURPLE
    s.line.width = Pt(2)
    stf2 = s.text_frame
    stf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    sr2 = stf2.paragraphs[0].add_run()
    sr2.text = f"{name}\n{desc}"
    sr2.font.size = Pt(14)
    sr2.font.color.rgb = WHITE

# Bullets at bottom
txBox2 = add_textbox(slide10, Inches(0.8), Inches(5.4), Inches(11.5), Inches(1.8))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_bullet(tf2, "Each agent focuses on what it does best", font_size=18)
add_bullet(tf2, "Agents communicate through messages or shared state", font_size=18)
add_rich_bullet(tf2, [("Frameworks like ", LIGHT_TEXT, False, False), ("CrewAI", LIGHT_TEXT, True, False), (", ", LIGHT_TEXT, False, False), ("AutoGen", LIGHT_TEXT, True, False), (", and ", LIGHT_TEXT, False, False), ("LangGraph", LIGHT_TEXT, True, False), (" make this easier", LIGHT_TEXT, False, False)], font_size=18)

# ============================================================
# SLIDE 11: Real-World Applications
# ============================================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide11, DARK_BG2)
add_slide_title(slide11, "Real-World Applications of Agentic AI")

add_table(slide11,
    ["Domain", "Application"],
    [
        ["Software Engineering", "Coding agents that write, debug, and deploy code (e.g., Claude Code, Devin)"],
        ["Research", "Agents that read papers, summarize findings, and generate hypotheses"],
        ["Customer Support", "Autonomous agents that resolve tickets, look up accounts, and escalate issues"],
        ["Data Analysis", "Agents that query databases, create charts, and write reports"],
        ["Business Operations", "Workflow automation \u2014 scheduling, invoicing, email management"],
        ["Education", "Personalized tutoring agents that adapt to each student's pace"],
    ],
    col_widths=[Inches(3.0), Inches(8.5)],
)

# ============================================================
# SLIDE 12: Popular Agent Frameworks
# ============================================================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide12, DARK_BG2)
add_slide_title(slide12, "Popular Agent Frameworks")

add_table(slide12,
    ["Framework", "Best For"],
    [
        ["LangChain / LangGraph", "General-purpose agent building with flexible tool integration"],
        ["CrewAI", "Multi-agent collaboration with role-based agents"],
        ["AutoGen (Microsoft)", "Conversational multi-agent patterns and code execution"],
        ["OpenClaw", "Our platform \u2014 hands-on agent building and deployment"],
        ["Anthropic Agent SDK", "Building agents powered by Claude with native tool use"],
    ],
    col_widths=[Inches(3.5), Inches(8.0)],
)

add_highlight_box(slide12, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.9), [
    [("In this course, we will primarily use ", LIGHT_TEXT, False, False),
     ("OpenClaw", LIGHT_TEXT, True, False),
     (" along with LangChain and CrewAI for hands-on labs.", LIGHT_TEXT, False, False)]
])

# ============================================================
# SLIDE 13: Safety, Ethics & Guardrails
# ============================================================
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide13, DARK_BG2)
add_slide_title(slide13, "Safety, Ethics & Guardrails")

txBox = add_textbox(slide13, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7))
tf = txBox.text_frame
tf.word_wrap = True
add_rich_paragraph(tf, [
    ("Autonomous agents are powerful \u2014 but they need ", LIGHT_TEXT, False, False),
    ("boundaries", LIGHT_TEXT, True, False),
    (" to operate safely.", LIGHT_TEXT, False, False),
], font_size=20)

txBox2 = add_textbox(slide13, Inches(0.8), Inches(2.5), Inches(11.5), Inches(4.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_rich_bullet(tf2, [("Hallucination Risk: ", BLUE, True, False), ("Agents can act on incorrect information \u2014 always verify critical outputs", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf2, [("Infinite Loops: ", BLUE, True, False), ("Without limits, agents may repeat actions endlessly \u2014 set max iterations", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf2, [("Harmful Actions: ", BLUE, True, False), ("Agents with file/system access need sandboxing and permission controls", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf2, [("Human-in-the-Loop: ", BLUE, True, False), ("For high-stakes decisions, require human approval before execution", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf2, [("Transparency: ", BLUE, True, False), ("Log all agent reasoning and actions for auditability", LIGHT_TEXT, False, False)], font_size=18)
add_rich_bullet(tf2, [("Bias & Fairness: ", BLUE, True, False), ("Agents inherit biases from their training data \u2014 test for and mitigate bias", LIGHT_TEXT, False, False)], font_size=18)

# ============================================================
# SLIDE 14: Course Structure Overview
# ============================================================
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide14, DARK_BG2)
add_slide_title(slide14, "Course Structure Overview")

add_table(slide14,
    ["Week", "Topic"],
    [
        ["1", "Introduction to Agentic AI & The Agent Landscape"],
        ["2", "LLMs as Reasoning Engines & The ReAct Pattern"],
        ["3", "Tool Use & External Integrations"],
        ["4", "Memory & Context Management"],
        ["5", "Planning Architectures & Self-Reflection"],
        ["6", "Multi-Agent Systems"],
        ["7", "Safety, Ethics & Evaluation"],
        ["8", "Capstone Projects & Future Directions"],
    ],
    col_widths=[Inches(1.5), Inches(10.0)],
)

# ============================================================
# SLIDE 15: What You Will Build
# ============================================================
slide15 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide15, DARK_BG2)
add_slide_title(slide15, "What You Will Build")

txBox = add_textbox(slide15, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.0))
tf = txBox.text_frame
tf.word_wrap = True
add_rich_bullet(tf, [("Lab 1: ", BLUE, True, False), ("A ReAct agent from scratch \u2014 observe the think-act-observe loop in action", LIGHT_TEXT, False, False)], font_size=19)
add_rich_bullet(tf, [("Lab 2: ", BLUE, True, False), ("A tool-augmented Q&A agent \u2014 search the web, read files, answer questions", LIGHT_TEXT, False, False)], font_size=19)
add_rich_bullet(tf, [("Lab 3: ", BLUE, True, False), ("An agent with persistent memory \u2014 remember past conversations using vector search", LIGHT_TEXT, False, False)], font_size=19)
add_rich_bullet(tf, [("Lab 4: ", BLUE, True, False), ("A multi-agent research team \u2014 researcher, writer, and reviewer working together", LIGHT_TEXT, False, False)], font_size=19)
add_rich_bullet(tf, [("Capstone: ", BLUE, True, False), ("A complete agentic system solving a real-world problem of your choice", LIGHT_TEXT, False, False)], font_size=19)

add_highlight_box(slide15, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.0), [
    [("By the end of this course, you won't just understand Agentic AI \u2014 you'll have ", LIGHT_TEXT, False, False),
     ("built", LIGHT_TEXT, True, False),
     (" multiple working agents.", LIGHT_TEXT, False, False)]
])

# ============================================================
# SLIDE 16: End Slide
# ============================================================
slide16 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide16, RGBColor(0x1A, 0x1A, 0x4E))

txBox = add_textbox(slide16, Inches(1), Inches(2.5), Inches(11.3), Inches(1.2))
tf = txBox.text_frame
set_text(tf, "Let's Build the Future of AI", font_size=52, color=BLUE, bold=True,
         alignment=PP_ALIGN.CENTER)

txBox2 = add_textbox(slide16, Inches(1), Inches(3.8), Inches(11.3), Inches(0.8))
tf2 = txBox2.text_frame
set_text(tf2, "Welcome to Agentic AI \u2014 OLI-AAI-101", font_size=24, color=MUTED_TEXT,
         alignment=PP_ALIGN.CENTER)

txBox3 = add_textbox(slide16, Inches(1), Inches(4.6), Inches(11.3), Inches(0.6))
tf3 = txBox3.text_frame
set_text(tf3, "OpenClaw Learning Initiative", font_size=18, color=DIM_TEXT,
         alignment=PP_ALIGN.CENTER)


# Save
output_path = r"c:\Users\user1\Desktop\OLICourses\Agentic-AI-Course\Course-Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
