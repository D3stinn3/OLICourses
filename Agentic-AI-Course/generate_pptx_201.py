"""
Generate a PowerPoint presentation for the OLI-AAI-201 course.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors from the HTML theme
DARK_BG = RGBColor(0x0A, 0x0A, 0x1A)
DARK_BG2 = RGBColor(0x0F, 0x0F, 0x2E)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
PINK = RGBColor(0xF4, 0x72, 0xB6)
LIGHT_TEXT = RGBColor(0xCB, 0xD5, 0xE1)
MUTED_TEXT = RGBColor(0x94, 0xA3, 0xB8)
DIM_TEXT = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xE0, 0xE0, 0xE0)
HIGHLIGHT_BG = RGBColor(0x1E, 0x3A, 0x5F)
TABLE_HEADER_BG = RGBColor(0x15, 0x2A, 0x4A)
TABLE_ROW_BG = RGBColor(0x0F, 0x14, 0x2E)
TABLE_ALT_BG = RGBColor(0x12, 0x1A, 0x35)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, font_size=18, color=LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_paragraph(tf, text, font_size=18, color=LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                   space_before=Pt(0), space_after=Pt(6)):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_rich_paragraph(tf, segments, font_size=18, alignment=PP_ALIGN.LEFT,
                        space_before=Pt(0), space_after=Pt(6)):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    for text, color, bold, italic in segments:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
    return p


def add_bullet(tf, text, font_size=20, color=LIGHT_TEXT, bullet_color=PURPLE,
               space_before=Pt(4), space_after=Pt(6)):
    p = tf.add_paragraph()
    p.space_before = space_before
    p.space_after = space_after

    bullet_run = p.add_run()
    bullet_run.text = "\u25B8  "
    bullet_run.font.size = Pt(font_size)
    bullet_run.font.color.rgb = bullet_color
    bullet_run.font.bold = True

    text_run = p.add_run()
    text_run.text = text
    text_run.font.size = Pt(font_size)
    text_run.font.color.rgb = color
    return p


def add_rich_bullet(tf, segments, font_size=20, bullet_color=PURPLE,
                     space_before=Pt(4), space_after=Pt(6)):
    p = tf.add_paragraph()
    p.space_before = space_before
    p.space_after = space_after

    bullet_run = p.add_run()
    bullet_run.text = "\u25B8  "
    bullet_run.font.size = Pt(font_size)
    bullet_run.font.color.rgb = bullet_color
    bullet_run.font.bold = True

    for text, color, bold, italic in segments:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
    return p


def add_highlight_box(slide, left, top, width, height, text_segments, font_size=18):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x12, 0x25, 0x40)
    shape.line.fill.background()

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    txBox = add_textbox(slide, left + Inches(0.3), top + Inches(0.15),
                        width - Inches(0.5), height - Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for segments in text_segments:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        for text, color, bold, italic in segments:
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.italic = italic


def add_diagram_flow(slide, items, top, left_start=None, box_width=Inches(2.2),
                      box_height=Inches(1.0), gap=Inches(0.3), arrow_width=Inches(0.5)):
    total_width = len(items) * box_width + (len(items) - 1) * (arrow_width + gap * 2)
    if left_start is None:
        left_start = (SLIDE_W - total_width) // 2

    x = left_start
    for i, (title, subtitle) in enumerate(items):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_width, box_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x1A, 0x15, 0x35)
        shape.line.color.rgb = PURPLE
        shape.line.width = Pt(2)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = title
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = WHITE

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(12)
        r2.font.color.rgb = MUTED_TEXT

        x += box_width + gap

        if i < len(items) - 1:
            arrow_tb = add_textbox(slide, x, top + Inches(0.25), arrow_width, Inches(0.5))
            atf = arrow_tb.text_frame
            atf.paragraphs[0].alignment = PP_ALIGN.CENTER
            ar = atf.paragraphs[0].add_run()
            ar.text = "\u2192"
            ar.font.size = Pt(28)
            ar.font.color.rgb = BLUE
            x += arrow_width + gap


def add_slide_title(slide, text, left=Inches(0.8), top=Inches(0.5), width=Inches(11.5)):
    txBox = add_textbox(slide, left, top, width, Inches(0.8))
    tf = txBox.text_frame
    set_text(tf, text, font_size=36, color=BLUE, bold=True)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top + Inches(0.75), Inches(10), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = HIGHLIGHT_BG
    line.line.fill.background()


def add_table(slide, headers, rows, left=Inches(0.8), top=Inches(2.0), width=Inches(11.5),
              col_widths=None):
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, Inches(0.5 * num_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(16)
                run.font.color.rgb = BLUE
                run.font.bold = True

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            bg = TABLE_ROW_BG if r % 2 == 0 else TABLE_ALT_BG
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(14)
                    run.font.color.rgb = LIGHT_TEXT

    tbl = table._tbl
    tblPr = tbl.tblPr if tbl.tblPr is not None else tbl._add_tblPr()
    tblPr.set('firstRow', '1')


# Column layout constants
left_x = Inches(0.8)
right_x = Inches(7.0)
col_w = Inches(5.5)

# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, RGBColor(0x0A, 0x0A, 0x2E))

txBox = add_textbox(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(1.5))
tf = txBox.text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
run = tf.paragraphs[0].add_run()
run.text = "Agentic AI: Agents Building in the Open"
run.font.size = Pt(52)
run.font.bold = True
run.font.color.rgb = BLUE

txBox2 = add_textbox(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.8))
tf2 = txBox2.text_frame
set_text(tf2, "Learn by Building the Platform You\u2019re Learning On", font_size=28, color=MUTED_TEXT,
         alignment=PP_ALIGN.CENTER)

txBox3 = add_textbox(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.6))
tf3 = txBox3.text_frame
set_text(tf3, "OLI-AAI-201 | OpenClaw Learning Initiative", font_size=16, color=DIM_TEXT,
         alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Welcome to OLI-AAI-201
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG2)
add_slide_title(slide, "Welcome to OLI-AAI-201")

txBox = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
add_rich_paragraph(tf, [
    ("In OLI-AAI-101 you learned ", LIGHT_TEXT, False, False),
    ("what", LIGHT_TEXT, True, False),
    (" AI agents are. In this course, you\u2019ll learn how agents ", LIGHT_TEXT, False, False),
    ("build real software", LIGHT_TEXT, True, False),
    (" \u2014 by contributing to the very platform you\u2019re using right now.", LIGHT_TEXT, False, False),
], font_size=20)

add_highlight_box(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(1.1), [
    [("The Meta-Learning Model: ", BLUE, True, False),
     ("This course was partially generated by AI agents. The slides you\u2019re reading, the quizzes you\u2019ll take, and the features you\u2019ll build all live in the same open-source codebase.", LIGHT_TEXT, False, False)]
])

txBox2 = add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(2.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_rich_bullet(tf2, [("You are both a ", LIGHT_TEXT, False, False), ("student", LIGHT_TEXT, True, False), (" and an ", LIGHT_TEXT, False, False), ("open-source contributor", LIGHT_TEXT, True, False)], font_size=20)
add_rich_bullet(tf2, [("AI agents are your ", LIGHT_TEXT, False, False), ("pair programming partners", LIGHT_TEXT, True, False), (", not replacements", LIGHT_TEXT, False, False)], font_size=20)
add_bullet(tf2, "Every feature you ship becomes part of the platform for future students", font_size=20)

# ============================================================
# SLIDE 3: The Scwripts Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG2)
add_slide_title(slide, "The Scwripts Architecture")

txBox = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7))
tf = txBox.text_frame
tf.word_wrap = True
add_paragraph(tf, "Scwripts is a full-stack education platform. Understanding its architecture is the first step to contributing.", font_size=20, color=LIGHT_TEXT)

# Left column: Backend
txL = add_textbox(slide, left_x, Inches(2.5), col_w, Inches(0.5))
set_text(txL.text_frame, "Backend (Django + Ninja)", font_size=24, color=PURPLE, bold=True)

txL2 = add_textbox(slide, left_x, Inches(3.1), col_w, Inches(3.5))
tfL = txL2.text_frame
tfL.word_wrap = True
for item in [
    "Django 6 with Django Ninja Extra API",
    "Controller \u2192 Service \u2192 Model pattern",
    "JWT authentication (ninja-jwt)",
    "Anthropic Claude for AI tutoring",
    "SQLite database (dev), PostgreSQL (prod)",
]:
    add_bullet(tfL, item, font_size=17)

# Right column: Frontend
txR = add_textbox(slide, right_x, Inches(2.5), col_w, Inches(0.5))
set_text(txR.text_frame, "Frontend (Next.js + React)", font_size=24, color=PURPLE, bold=True)

txR2 = add_textbox(slide, right_x, Inches(3.1), col_w, Inches(3.5))
tfR = txR2.text_frame
tfR.word_wrap = True
for item in [
    "Next.js Pages Router with dynamic routes",
    "React Context for auth and state",
    "React Three Fiber for 3D visualizations",
    "CSS Modules for scoped styling",
    "Centralized API client with auto token refresh",
]:
    add_bullet(tfR, item, font_size=17)

# ============================================================
# SLIDE 4: The Six Backend Apps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG2)
add_slide_title(slide, "The Six Backend Apps")

txBox = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5))
tf = txBox.text_frame
tf.word_wrap = True
add_rich_paragraph(tf, [
    ("The Scwripts backend is organized into ", LIGHT_TEXT, False, False),
    ("six Django apps", LIGHT_TEXT, True, False),
    (", each owning a specific domain.", LIGHT_TEXT, False, False),
], font_size=20)

add_table(slide,
    ["App", "Purpose", "Key Models"],
    [
        ["accounts", "User registration and JWT auth", "User (Django built-in)"],
        ["courses", "Course catalog, modules, slides, enrollment", "Course, Module, Slide, Enrollment"],
        ["quizzes", "Assessments and grading", "Quiz, Question, Choice, Attempt"],
        ["agentic", "AI tutor chat with streaming", "ChatSession, ChatMessage"],
        ["gamification", "XP, levels, achievements, leagues, quests", "UserXP, Achievement, DailyQuest, UserAgent"],
        ["engagement", "Facial analysis and adaptive learning", "EngagementSnapshot, EngagementSummary, EngagementConsent"],
    ],
    col_widths=[Inches(2.0), Inches(5.0), Inches(4.5)],
)

# ============================================================
# SLIDE 5: How Agents Built This Platform
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG2)
add_slide_title(slide, "How Agents Built This Platform")

txBox = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5))
tf = txBox.text_frame
tf.word_wrap = True
add_paragraph(tf, "Every major feature in Scwripts was built with AI agent assistance. Let\u2019s trace how agents contributed to the codebase.", font_size=20, color=LIGHT_TEXT)

txBox2 = add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(3.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_rich_bullet(tf2, [("Course Content Generation: ", BLUE, True, False), ("An agent produced the entire AAI-101 slide deck and quiz as structured JSON \u2014 stored in seed_agentic_course.py", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Backend Architecture: ", BLUE, True, False), ("Agents scaffolded the Controller \u2192 Service \u2192 Model pattern across all six apps", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("AI Tutor Service: ", BLUE, True, False), ("The LLMService streaming implementation was pair-programmed with an agent using Anthropic\u2019s SDK", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Gamification System: ", BLUE, True, False), ("XP thresholds, level titles, and achievement criteria were designed and coded by an agent", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("3D Visualizations: ", BLUE, True, False), ("React Three Fiber scenes (ReAct loop, Multi-Agent System, Planning Tree) were agent-generated", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("This Course: ", BLUE, True, False), ("The slides you\u2019re reading right now were generated by an agent and reviewed by a human instructor", LIGHT_TEXT, False, False)], font_size=17)

add_highlight_box(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0), [
    [("Transparency: ", BLUE, True, False),
     ("In this course, you\u2019ll always know what the agent built vs. what a human built. This is a core principle of responsible AI-augmented development.", LIGHT_TEXT, False, False)]
])

# ============================================================
# SLIDE 6: Inside the Slide Engine
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG2)
add_slide_title(slide, "Inside the Slide Engine")

txBox = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5))
tf = txBox.text_frame
tf.word_wrap = True
add_paragraph(tf, "Let\u2019s understand exactly how course content flows from agent generation to what you see on screen.", font_size=20, color=LIGHT_TEXT)

add_diagram_flow(slide, [
    ("Agent", "Generates JSON"),
    ("Seed Command", "Loads into DB"),
    ("API", "Serves slides"),
    ("SlideContent", "Renders UI"),
], top=Inches(2.3))

txBox2 = add_textbox(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(3.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_rich_bullet(tf2, [("Step 1 \u2014 Generation: ", BLUE, True, False), ("An agent produces an array of slide objects with type, heading, paragraphs, bullets, diagram, columns, tableHeaders/Rows, highlightBox", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Step 2 \u2014 Seeding: ", BLUE, True, False), ("A Django management command (python manage.py seed_*) creates Course, Module, and Slide records in the database", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Step 3 \u2014 API: ", BLUE, True, False), ("The CourseController serves slides as JSON via /api/courses/{slug}/slides", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Step 4 \u2014 Rendering: ", BLUE, True, False), ("The SlideContent React component reads the JSON and renders the appropriate layout for each slide type", LIGHT_TEXT, False, False)], font_size=17)

# ============================================================
# SLIDE 7: The Slide JSON Schema
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG2)
add_slide_title(slide, "The Slide JSON Schema")

txBox = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7))
tf = txBox.text_frame
tf.word_wrap = True
add_rich_paragraph(tf, [
    ("When an agent generates course slides, it must produce JSON that matches this schema exactly. This is ", LIGHT_TEXT, False, False),
    ("structured output", LIGHT_TEXT, True, False),
    (" in action.", LIGHT_TEXT, False, False),
], font_size=20)

# Left column: Title/End slides
txL = add_textbox(slide, left_x, Inches(2.5), col_w, Inches(0.5))
set_text(txL.text_frame, "Title / End Slides", font_size=24, color=PURPLE, bold=True)

txL2 = add_textbox(slide, left_x, Inches(3.1), col_w, Inches(2.5))
tfL = txL2.text_frame
tfL.word_wrap = True
for item in [
    "type: 'title' or 'end'",
    "title: Main heading text",
    "subtitle: Secondary text",
    "meta: Course code and program",
]:
    add_bullet(tfL, item, font_size=17)

# Right column: Content slides
txR = add_textbox(slide, right_x, Inches(2.5), col_w, Inches(0.5))
set_text(txR.text_frame, "Content Slides", font_size=24, color=PURPLE, bold=True)

txR2 = add_textbox(slide, right_x, Inches(3.1), col_w, Inches(3.5))
tfR = txR2.text_frame
tfR.word_wrap = True
for item in [
    "type: 'content'",
    "heading: Slide heading (required)",
    "paragraphs: Array of HTML strings",
    "bullets: Array of HTML strings",
    "diagram: Array of {label, sub, accent?}",
    "columns: Array of {title, items[]}",
    "tableHeaders + tableRows: Table data",
    "highlightBox: Callout HTML string",
]:
    add_bullet(tfR, item, font_size=17)

add_highlight_box(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0), [
    [("Lab 1 Challenge: ", BLUE, True, False),
     ("You\u2019ll build an agent that generates valid slide JSON, validates it against this schema, and loads it into the database via a seed command.", LIGHT_TEXT, False, False)]
])

# ============================================================
# SLIDE 8: The Streaming Chat Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG2)
add_slide_title(slide, "The Streaming Chat Architecture")

txBox = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5))
tf = txBox.text_frame
tf.word_wrap = True
add_paragraph(tf, "The AI Tutor is one of Scwripts\u2019 most sophisticated features. Let\u2019s trace how a student message becomes a streamed AI response.", font_size=20, color=LIGHT_TEXT)

add_diagram_flow(slide, [
    ("Student", "Sends message"),
    ("ChatController", "Authenticates + saves"),
    ("LLMService", "Calls Claude API"),
    ("SSE Stream", "Token-by-token"),
], top=Inches(2.3))

txBox2 = add_textbox(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(3.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_rich_bullet(tf2, [("Frontend: ", BLUE, True, False), ("streamChat() in api.js sends a POST to /api/chat/stream with the message history and slide context", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Controller: ", BLUE, True, False), ("Authenticates the user, loads or creates a ChatSession, saves the user message as a ChatMessage", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("LLMService: ", BLUE, True, False), ("Builds the system prompt (including slide context and engagement data), calls anthropic.messages.stream()", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Response: ", BLUE, True, False), ("Each text token is yielded as data: {token}\\n\\n (Server-Sent Events format), displayed in real-time on the frontend", LIGHT_TEXT, False, False)], font_size=17)

# ============================================================
# SLIDE 9: The Gamification Engine
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG2)
add_slide_title(slide, "The Gamification Engine")

txBox = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5))
tf = txBox.text_frame
tf.word_wrap = True
add_paragraph(tf, "Scwripts uses a comprehensive gamification system to keep students motivated. This entire system was designed and implemented with agent assistance.", font_size=20, color=LIGHT_TEXT)

add_table(slide,
    ["Feature", "How It Works", "Agent Enhancement Opportunity"],
    [
        ["XP & Levels", "Users earn XP from lessons, quizzes, chat, challenges. Levels 1-50 with titles.", "Agents could adjust XP rewards based on difficulty"],
        ["Achievements", "Badge-like rewards with criteria (e.g., 'Complete 5 quizzes')", "Agents could generate new achievements based on usage patterns"],
        ["Daily Quests", "Time-limited challenges refreshed daily", "Agents could personalize quests based on student progress"],
        ["Leagues", "Weekly leaderboards with tiers (Reactive \u2192 Superintelligence)", "Agents could balance matchmaking for fair competition"],
        ["Agent Companion", "Virtual agent that levels up with the student", "Students build capabilities for their companion agent in labs"],
    ],
    col_widths=[Inches(2.0), Inches(5.0), Inches(4.5)],
)

# ============================================================
# SLIDE 10: The Engagement System
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG2)
add_slide_title(slide, "The Engagement System")

txBox = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7))
tf = txBox.text_frame
tf.word_wrap = True
add_rich_paragraph(tf, [
    ("Scwripts includes an ", LIGHT_TEXT, False, False),
    ("opt-in", LIGHT_TEXT, True, False),
    (" engagement tracking system that uses facial analysis to measure student attention and adapt content delivery.", LIGHT_TEXT, False, False),
], font_size=20)

txBox2 = add_textbox(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(3.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_rich_bullet(tf2, [("Consent First: ", BLUE, True, False), ("The EngagementConsent model ensures students explicitly opt in before any tracking begins", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Facial Analysis: ", BLUE, True, False), ("The FaceService processes webcam snapshots to detect emotion, gaze direction, and overall engagement", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Snapshots: ", BLUE, True, False), ("Each analysis creates an EngagementSnapshot with emotion, confidence, gaze pitch/yaw, and engagement score", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Daily Summaries: ", BLUE, True, False), ("The system aggregates snapshots into EngagementSummary records with average engagement, dominant emotion, and attention drops", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Adaptive Response: ", BLUE, True, False), ("The AdaptiveService adjusts AI tutor behavior based on engagement \u2014 e.g., simplifying explanations when confusion is detected", LIGHT_TEXT, False, False)], font_size=17)

add_highlight_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.0), [
    [("Privacy by Design: ", BLUE, True, False),
     ("Raw images are never stored. Only derived metrics are saved. Students can revoke consent at any time, and all their engagement data is deleted.", LIGHT_TEXT, False, False)]
])

# ============================================================
# SLIDE 11: Open-Source Contribution Workflow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG2)
add_slide_title(slide, "Open-Source Contribution Workflow")

txBox = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5))
tf = txBox.text_frame
tf.word_wrap = True
add_paragraph(tf, "In this course, you contribute to Scwripts like any open-source project. Here\u2019s the workflow you\u2019ll follow for every lab and assignment.", font_size=20, color=LIGHT_TEXT)

add_diagram_flow(slide, [
    ("Fork", "Copy the repo"),
    ("Branch", "Create feature branch"),
    ("Build", "Code with agent help"),
    ("PR", "Submit for review"),
], top=Inches(2.3))

txBox2 = add_textbox(slide, Inches(0.8), Inches(3.6), Inches(11.5), Inches(3.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_rich_bullet(tf2, [("1. Fork ", BLUE, True, False), ("the Scwripts repository to your GitHub account", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("2. Branch ", BLUE, True, False), ("from main with a descriptive name (e.g., feat/study-guide-api)", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("3. Build ", BLUE, True, False), ("your feature using an AI agent as a pair programmer \u2014 document what the agent did vs. what you did", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("4. Test ", BLUE, True, False), ("your changes locally (backend + frontend) and write automated tests", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("5. Submit a PR ", BLUE, True, False), ("with a clear description, screenshots/demo, and the agent collaboration log", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("6. Review ", BLUE, True, False), ("at least one peer\u2019s PR and incorporate feedback on your own", LIGHT_TEXT, False, False)], font_size=17)

# ============================================================
# SLIDE 12: Building a Content Generation Agent
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG2)
add_slide_title(slide, "Building a Content Generation Agent")

txBox = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5))
tf = txBox.text_frame
tf.word_wrap = True
add_paragraph(tf, "Your first major lab is to build an agent that generates course content. Here\u2019s the architecture of a content generation pipeline.", font_size=20, color=LIGHT_TEXT)

add_diagram_flow(slide, [
    ("Topic Input", "User provides subject"),
    ("Research Agent", "Gathers key concepts"),
    ("Writer Agent", "Generates slide JSON"),
    ("Validator", "Checks schema + quality"),
], top=Inches(2.3))

txBox2 = add_textbox(slide, Inches(0.8), Inches(3.6), Inches(11.5), Inches(2.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_rich_bullet(tf2, [("Input: ", BLUE, True, False), ("A topic string and target audience level (beginner, intermediate, advanced)", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Research Phase: ", BLUE, True, False), ("The agent identifies 8-15 key concepts, definitions, and examples for the topic", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Generation Phase: ", BLUE, True, False), ("The agent produces slide JSON matching the Scwripts schema \u2014 using structured output / JSON mode", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Validation Phase: ", BLUE, True, False), ("A separate pass checks that all required fields are present, HTML is valid, and content is accurate", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Output: ", BLUE, True, False), ("A Python seed command file that can be run with python manage.py seed_[topic]", LIGHT_TEXT, False, False)], font_size=17)

add_highlight_box(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.9), [
    [("Lab 1: ", BLUE, True, False),
     ("Build this pipeline. Your agent must generate at least 10 valid content slides and 8 quiz questions for any given topic.", LIGHT_TEXT, False, False)]
])

# ============================================================
# SLIDE 13: Agent Pair Programming: Backend Features
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG2)
add_slide_title(slide, "Agent Pair Programming: Backend Features")

txBox = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5))
tf = txBox.text_frame
tf.word_wrap = True
add_paragraph(tf, "In Lab 2, you\u2019ll implement a complete backend feature with an AI agent as your pair programmer. Here\u2019s how the collaboration works.", font_size=20, color=LIGHT_TEXT)

# Left column: What the Agent Does
txL = add_textbox(slide, left_x, Inches(2.3), col_w, Inches(0.5))
set_text(txL.text_frame, "What the Agent Does", font_size=24, color=PURPLE, bold=True)

txL2 = add_textbox(slide, left_x, Inches(2.9), col_w, Inches(3.0))
tfL = txL2.text_frame
tfL.word_wrap = True
for item in [
    "Scaffolds model, schema, controller, and service files",
    "Generates database migration code",
    "Writes initial test cases",
    "Suggests API endpoint designs",
    "Handles boilerplate and repetitive patterns",
]:
    add_bullet(tfL, item, font_size=17)

# Right column: What You Do
txR = add_textbox(slide, right_x, Inches(2.3), col_w, Inches(0.5))
set_text(txR.text_frame, "What You Do", font_size=24, color=PURPLE, bold=True)

txR2 = add_textbox(slide, right_x, Inches(2.9), col_w, Inches(3.0))
tfR = txR2.text_frame
tfR.word_wrap = True
for item in [
    "Define the feature requirements and data model",
    "Review and correct agent-generated code",
    "Make architectural decisions (where logic lives)",
    "Write edge case tests the agent missed",
    "Ensure security, privacy, and error handling",
]:
    add_bullet(tfR, item, font_size=17)

add_highlight_box(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0), [
    [("Lab 2: ", BLUE, True, False),
     ("Implement a StudyGuide feature \u2014 model, schema, controller, service, migration, and tests. Document every agent interaction.", LIGHT_TEXT, False, False)]
])

# ============================================================
# SLIDE 14: Agent Pair Programming: Frontend Features
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG2)
add_slide_title(slide, "Agent Pair Programming: Frontend Features")

txBox = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5))
tf = txBox.text_frame
tf.word_wrap = True
add_paragraph(tf, "In Lab 3, you\u2019ll build a complete frontend feature. The agent helps with component structure, styling, and API integration.", font_size=20, color=LIGHT_TEXT)

txBox2 = add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(3.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_rich_bullet(tf2, [("Page Design: ", BLUE, True, False), ("Tell the agent what the page should do \u2014 it generates a React component with the right imports, hooks, and layout", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("API Integration: ", BLUE, True, False), ("The agent wires up calls to api.js with proper auth token handling and error states", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Styling: ", BLUE, True, False), ("The agent creates a CSS Module matching the existing Scwripts design system (dark theme, consistent spacing, responsive layout)", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("State Management: ", BLUE, True, False), ("The agent suggests when to use local state vs. Context vs. a new provider", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Your Role: ", BLUE, True, False), ("Define UX requirements, review accessibility, test edge cases (loading, error, empty states), and ensure mobile responsiveness", LIGHT_TEXT, False, False)], font_size=17)

add_highlight_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.2), [
    [("Lab 3: ", BLUE, True, False),
     ("Build an \"Agent Lab\" page \u2014 an interactive sandbox where students can test agent prompts, see tool calls visualized, and compare different agent strategies.", LIGHT_TEXT, False, False)]
])

# ============================================================
# SLIDE 15: Multi-Agent Content Pipeline
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG2)
add_slide_title(slide, "Multi-Agent Content Pipeline")

txBox = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5))
tf = txBox.text_frame
tf.word_wrap = True
add_rich_paragraph(tf, [
    ("The capstone of this course is building a ", LIGHT_TEXT, False, False),
    ("multi-agent pipeline", LIGHT_TEXT, True, False),
    (" that can generate an entire course module autonomously.", LIGHT_TEXT, False, False),
], font_size=20)

add_diagram_flow(slide, [
    ("Orchestrator", "Manages workflow"),
    ("Researcher", "Gathers concepts"),
    ("Writer", "Generates slides"),
    ("Reviewer", "Validates quality"),
], top=Inches(2.3))

txBox2 = add_textbox(slide, Inches(0.8), Inches(3.6), Inches(11.5), Inches(2.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_rich_bullet(tf2, [("Orchestrator Agent: ", BLUE, True, False), ("Receives a topic, breaks it into subtasks, delegates to specialists, and manages the review loop", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Research Agent: ", BLUE, True, False), ("Searches for key concepts, definitions, real-world examples, and common misconceptions", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Writer Agent: ", BLUE, True, False), ("Produces slide JSON and quiz questions in the Scwripts schema format", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Reviewer Agent: ", BLUE, True, False), ("Validates content accuracy, checks schema compliance, scores quality, and requests revisions if needed", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Quiz Agent: ", BLUE, True, False), ("Generates assessment questions with explanations, ensuring coverage of all key concepts", LIGHT_TEXT, False, False)], font_size=17)

add_highlight_box(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(1.0), [
    [("Lab 4: ", BLUE, True, False),
     ("Build this multi-agent pipeline. Given a single topic prompt, it must produce a complete module: 10+ slides, 8+ quiz questions, and a study guide \u2014 all schema-valid and quality-reviewed.", LIGHT_TEXT, False, False)]
])

# ============================================================
# SLIDE 16: Testing Agent-Generated Output
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG2)
add_slide_title(slide, "Testing Agent-Generated Output")

txBox = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5))
tf = txBox.text_frame
tf.word_wrap = True
add_paragraph(tf, "Agent outputs are non-deterministic. You can\u2019t test them with exact string matching. Here\u2019s how to test effectively.", font_size=20, color=LIGHT_TEXT)

add_table(slide,
    ["Testing Strategy", "What It Validates", "Example"],
    [
        ["Schema Validation", "Output structure matches expected format", "Every slide has a 'type' and 'content' field"],
        ["Constraint Checking", "Output meets defined boundaries", "Quiz has 4 choices per question, exactly 1 correct"],
        ["Quality Scoring", "Content meets a minimum quality bar", "Slide has >= 3 bullets, heading <= 60 chars"],
        ["Regression Testing", "New outputs don't break existing features", "Generated slides render correctly in SlideContent"],
        ["Human-in-the-Loop", "Final review before production deployment", "Instructor reviews and approves generated content"],
    ],
    col_widths=[Inches(2.5), Inches(4.5), Inches(4.5)],
)

# ============================================================
# SLIDE 17: Agent-in-the-Loop DevOps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG2)
add_slide_title(slide, "Agent-in-the-Loop DevOps")

txBox = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5))
tf = txBox.text_frame
tf.word_wrap = True
add_paragraph(tf, "Beyond content generation, agents can participate in the entire software development lifecycle.", font_size=20, color=LIGHT_TEXT)

txBox2 = add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(3.3))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_rich_bullet(tf2, [("Issue Triage: ", BLUE, True, False), ("An agent reads new GitHub issues, categorizes them (bug/feature/docs), estimates complexity, and suggests assignees", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("PR Generation: ", BLUE, True, False), ("Given an issue description, an agent creates a branch, implements the change, writes tests, and opens a PR", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Code Review: ", BLUE, True, False), ("An agent reviews PRs for correctness, security vulnerabilities, style consistency, and test coverage", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("CI/CD Integration: ", BLUE, True, False), ("Agent-powered checks run alongside traditional CI \u2014 e.g., verifying that generated content meets quality thresholds", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Documentation: ", BLUE, True, False), ("Agents auto-generate API docs, README updates, and changelog entries from commit history", LIGHT_TEXT, False, False)], font_size=17)

add_highlight_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.2), [
    [("In a mature agent-augmented workflow, humans focus on ", LIGHT_TEXT, False, False),
     ("design decisions, edge cases, and quality judgment", LIGHT_TEXT, True, False),
     (" while agents handle ", LIGHT_TEXT, False, False),
     ("scaffolding, boilerplate, and routine checks", LIGHT_TEXT, True, False),
     (".", LIGHT_TEXT, False, False)]
])

# ============================================================
# SLIDE 18: The Self-Improving Platform
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG2)
add_slide_title(slide, "The Self-Improving Platform")

txBox = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7))
tf = txBox.text_frame
tf.word_wrap = True
add_rich_paragraph(tf, [
    ("The ultimate vision for Scwripts is a ", LIGHT_TEXT, False, False),
    ("self-improving platform", LIGHT_TEXT, True, False),
    (" \u2014 where agents learn from student engagement data to generate better content over time.", LIGHT_TEXT, False, False),
], font_size=20)

add_diagram_flow(slide, [
    ("Students Learn", "Engagement data collected"),
    ("Agents Analyze", "Identify weak spots"),
    ("Agents Generate", "Better content"),
    ("Students Learn", "Improved outcomes"),
], top=Inches(2.5))

txBox2 = add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(3.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_rich_bullet(tf2, [("Feedback Loop: ", BLUE, True, False), ("Quiz scores, engagement snapshots, and chat patterns reveal which content works and which doesn\u2019t", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Adaptive Generation: ", BLUE, True, False), ("Agents regenerate slides that have low engagement or high confusion scores", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Difficulty Calibration: ", BLUE, True, False), ("Quiz questions are automatically adjusted based on pass/fail rates across cohorts", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Personalization: ", BLUE, True, False), ("The AI tutor adapts its teaching style based on individual student engagement patterns", LIGHT_TEXT, False, False)], font_size=17)
add_rich_bullet(tf2, [("Recursive Growth: ", BLUE, True, False), ("Each cohort\u2019s contributions make the platform better for the next cohort \u2014 a true open-source flywheel", LIGHT_TEXT, False, False)], font_size=17)

# ============================================================
# SLIDE 19: Course Structure Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG2)
add_slide_title(slide, "Course Structure Overview")

add_table(slide,
    ["Week", "Topic", "Deliverable"],
    [
        ["1", "The Codebase is the Classroom", "Lab 0: First PR"],
        ["2", "How Agents Generate Course Content", "Lab 1: Content Generation Agent"],
        ["3", "Agent-Powered Backend Development", "Lab 2: Backend Feature PR"],
        ["4", "Agent-Powered Frontend Development", "Lab 3: Frontend Feature PR"],
        ["5", "The Gamification Engine & Personalization", "Midterm Feature PR"],
        ["6", "Multi-Agent Content Pipelines", "Lab 4: Multi-Agent Pipeline"],
        ["7", "Agent-in-the-Loop DevOps & Testing", "Test Harness PR"],
        ["8", "Capstone & The Self-Improving Platform", "Capstone PR + Presentation"],
    ],
    col_widths=[Inches(1.5), Inches(5.5), Inches(4.5)],
)

# ============================================================
# SLIDE 20: What You Will Build
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG2)
add_slide_title(slide, "What You Will Build")

txBox = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.0))
tf = txBox.text_frame
tf.word_wrap = True
add_rich_bullet(tf, [("Lab 0: ", BLUE, True, False), ("Fork the Scwripts repo, run it locally, and submit your first pull request", LIGHT_TEXT, False, False)], font_size=19)
add_rich_bullet(tf, [("Lab 1: ", BLUE, True, False), ("A content generation agent that produces valid slide JSON and quiz questions from any topic", LIGHT_TEXT, False, False)], font_size=19)
add_rich_bullet(tf, [("Lab 2: ", BLUE, True, False), ("A complete backend feature (model + schema + controller + service + tests) built with agent pair programming", LIGHT_TEXT, False, False)], font_size=19)
add_rich_bullet(tf, [("Lab 3: ", BLUE, True, False), ("A complete frontend feature (page + components + API integration + styles) built with agent pair programming", LIGHT_TEXT, False, False)], font_size=19)
add_rich_bullet(tf, [("Lab 4: ", BLUE, True, False), ("A multi-agent pipeline that generates a full course module with research, writing, review, and assessment", LIGHT_TEXT, False, False)], font_size=19)
add_rich_bullet(tf, [("Capstone: ", BLUE, True, False), ("A significant feature contribution to Scwripts \u2014 shipped, tested, reviewed, and merged", LIGHT_TEXT, False, False)], font_size=19)

add_highlight_box(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.2), [
    [("By the end of this course, you won\u2019t just understand how agents build software \u2014 you\u2019ll have ", LIGHT_TEXT, False, False),
     ("shipped real features", LIGHT_TEXT, True, False),
     (" to a platform used by other students, with agents as your collaborators.", LIGHT_TEXT, False, False)]
])

# ============================================================
# SLIDE 21: End Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, RGBColor(0x1A, 0x1A, 0x4E))

txBox = add_textbox(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.2))
tf = txBox.text_frame
set_text(tf, "Let\u2019s Build in the Open", font_size=52, color=BLUE, bold=True,
         alignment=PP_ALIGN.CENTER)

txBox2 = add_textbox(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(0.8))
tf2 = txBox2.text_frame
set_text(tf2, "Welcome to Agentic AI: Agents Building in the Open \u2014 OLI-AAI-201", font_size=24, color=MUTED_TEXT,
         alignment=PP_ALIGN.CENTER)

txBox3 = add_textbox(slide, Inches(1), Inches(4.6), Inches(11.3), Inches(0.6))
tf3 = txBox3.text_frame
set_text(tf3, "OpenClaw Learning Initiative", font_size=18, color=DIM_TEXT,
         alignment=PP_ALIGN.CENTER)


# Save
output_path = r"c:\Users\user1\Desktop\OLICourses\Agentic-AI-Course\Course-Presentation-201.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
